from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import authenticate, login
from django.contrib.auth.hashers import make_password
from django.contrib import messages
from .models import *
from django.contrib.auth import get_user_model
from django.db.models import Q
from django.contrib.auth import logout
from django.shortcuts import redirect
from django.utils import timezone
import razorpay
from django.conf import settings
import os
from django.http import HttpResponse
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import google.generativeai as genai
from django.shortcuts import render
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from PyPDF2 import PdfReader
from pptx import Presentation
from .models import ChatHistory

#--------------------Chat_Model-----------------------------#
# Configure Gemini AI
genai.configure(api_key="")  # Replace with your actual API key
model = genai.GenerativeModel('gemini-1.5-pro')

# Function to get AI response from Gemini API
def get_ai_response(message, file_context=""):
    combined_input = f"File Content:\n{file_context}\n\nUser Query:\n{message}"
    chat_session = model.start_chat()
    response = chat_session.send_message(combined_input)
    return response.text


# Use the get_user_model function to get the custom user model
User = get_user_model()

def index(request):
    sizes = [5, 6, 7, 8, 9, 10, 11, 12]
    products = Product.objects.all()
    return render(request, 'home.html',{'products':products,'sizes': sizes})

# Create your views here.
def home(request):
    products = Product.objects.all()
    return render(request, 'home.html',{'products':products})

# Admin Dashboard
def admin_dashboard(request):
    products = Product.objects.all()
    return render(request, 'admin/admin_dashboard.html',{'products':products})



def approve_user(request, user_id):

    user = get_object_or_404(User, id=user_id)

    # Check if the user is already approved
    if user.is_approved:
        messages.info(request, f"User {user.username} is already approved.")
    else:
        # Approve the user
        user.is_approved = True
        user.save()
        messages.success(request, f"User {user.username} has been approved.")
    
    return redirect('admin_dashboard')

def approve_supplier(request, supplier_id):
    # Ensure that the logged-in user is an admin
    supplier = get_object_or_404(Supplier, id=supplier_id)
    
    # Check if the supplier is already approved
    if supplier.user.is_approved:
        messages.info(request, f"Supplier {supplier.company_name} is already approved.")
    else:
        # Approve the supplier
        supplier.user.is_approved = True
        supplier.user.save()
        messages.success(request, f"Supplier {supplier.company_name} has been approved.")
    
    return redirect('admin_dashboard')

def view_all_users(request):
    # Fetch all users who are not superusers and whose role is not 'supplier'
    users = User.objects.filter(is_superuser=False).exclude(role='supplier')
    
    # Render the template and pass the filtered list of users
    return render(request, 'admin/view_all_users.html', {'users': users})


def view_all_suppliers(request):
    # Fetch all suppliers
    suppliers = Supplier.objects.all()
    
    # Render the template and pass the list of suppliers
    return render(request, 'admin/view_all_suppliers.html', {'suppliers': suppliers})

def update_user(request, user_id):
    user = get_object_or_404(User, id=user_id)

    if request.method == 'POST':
        new_username = request.POST.get('username')
        new_email = request.POST.get('email')
        new_role = request.POST.get('role')

        # Validate the inputs (basic validation for example)
        if new_username and new_email and new_role in ['admin', 'supplier', 'user']:
            user.username = new_username
            user.email = new_email
            user.role = new_role
            user.save()
            messages.success(request, f'User {user.username} updated successfully.')
        else:
            messages.error(request, 'Invalid input.')

        return redirect('view_all_users')  # Redirect to the page showing all users

    return render(request, 'admin/update_user.html', {'user': user})

def update_supplier(request, supplier_id):
    supplier = get_object_or_404(Supplier, id=supplier_id)

    if request.method == 'POST':
        new_company_name = request.POST.get('company_name')
        new_contact_info = request.POST.get('contact_info')

        if new_company_name:
            supplier.company_name = new_company_name
        if new_contact_info:
            supplier.contact_info = new_contact_info

        supplier.save()
        messages.success(request, f'Supplier {supplier.company_name} updated successfully.')

        return redirect('view_all_suppliers')  # Redirect to the page showing all suppliers

    return render(request, 'admin/update_supplier.html', {'supplier': supplier})


def delete_user(request, user_id):
    user = get_object_or_404(User, id=user_id)

    if request.method == 'POST':  # Ensure deletion is a POST request
        user.delete()
        messages.success(request, f'User {user.username} deleted successfully.')
        return redirect('view_all_users')  # Redirect to the page showing all users

    return render(request, 'admin/delete_user.html', {'user': user})

def delete_supplier(request, supplier_id):
    supplier = get_object_or_404(Supplier, id=supplier_id)

    if request.method == 'POST':  # Ensure deletion is a POST request
        supplier.delete()
        messages.success(request, f'Supplier {supplier.company_name} deleted successfully.')
        return redirect('view_all_suppliers')  # Redirect to the page showing all suppliers

    return render(request, 'admin/delete_supplier.html', {'supplier': supplier})

#______________________________________________________________________________________________________________#

# Supplier Dashboard
def supplier_dashboard(request):
    return render(request, 'supplier/supplier_dashboard.html')

#______________________________________________________________________________________________________________#

# User Dashboard
def user_dashboard(request):
    products = Product.objects.all()
    sizes = [5, 6, 7, 8, 9, 10, 11, 12] 
    latest_orders = Order.objects.order_by('-created_at')[:10]
    trending_products = set()
    
    for order in latest_orders:
        trending_products.update(order.products.all()) # Define size options here
    return render(request, 'user/user_dashboard.html',{'products': products, 'sizes': sizes,'trending_products':trending_products})

def view_profile(request):
    # Get the current user's details
    user = request.user

    # Render the profile page with user details
    return render(request, 'user/view_profile.html', {'user': user})

def supplier_profile(request):
    # Get the current user's details
    user = request.user

    # Render the profile page with user details
    return render(request, 'supplier/view_profile.html', {'user': user})

def update_supplier_profile(request):
    # Get the current user's details
    user = request.user

    if request.method == 'POST':
        # Get the updated data from the POST request
        username = request.POST.get('username')
        email = request.POST.get('email')
        # Update user attributes
        user.username = username
        user.email = email
        user.save()

        # Show a success message
        messages.success(request, "Profile updated successfully!")
        return redirect('update_profile')  # Redirect to the same page after update

    return render(request, 'supplier/update_profile.html', {'user': user})


def update_profile(request):
    # Get the current user's details
    user = request.user

    if request.method == 'POST':
        # Get the updated data from the POST request
        username = request.POST.get('username')
        email = request.POST.get('email')
        # Update user attributes
        user.username = username
        user.email = email
        user.save()

        # Show a success message
        messages.success(request, "Profile updated successfully!")
        return redirect('update_profile')  # Redirect to the same page after update

    return render(request, 'user/update_profile.html', {'user': user})
#________________________________________________________________________________________#

#Chat_Bot(View)
# Function to extract text from uploaded PDF files
def extract_text_from_pdf(file):
    reader = PdfReader(file)
    return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

# Function to extract text from uploaded TXT files
def extract_text_from_txt(file):
    return file.read().decode('utf-8')

# Function to extract text from uploaded PPTX files
def extract_text_from_ppt(file):
    prs = Presentation(file)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

# Process uploaded file and extract text
def process_uploaded_file(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file)
    elif ext == ".txt":
        return extract_text_from_txt(file)
    elif ext == ".pptx":
        return extract_text_from_ppt(file)
    else:
        return "Unsupported file type."

# Chat page view
def chat_view(request):
    return render(request, 'user/chat.html')

# Handle chat response and store chat history
@csrf_exempt
def chat_response(request):
    if request.method == 'POST':
        message = request.POST.get('message')
        uploaded_file = request.FILES.get('uploadedFile')
        file_context = process_uploaded_file(uploaded_file) if uploaded_file else ""

        # Get AI response
        ai_response = get_ai_response(message, file_context)

        # Save chat history
        if request.user.is_authenticated:
            ChatHistory.objects.create(
                user=request.user,
                user_message=message,
                ai_response=ai_response
            )

        return JsonResponse({'response': ai_response})
    
    return JsonResponse({'error': 'Invalid request'}, status=400)

# Fetch user chat history
def get_chat_history(request):
    if request.user.is_authenticated:
        chats = ChatHistory.objects.filter(user=request.user).values("user_message", "ai_response", "timestamp")
        return JsonResponse({"chat_history": list(chats)})
    return JsonResponse({"error": "Unauthorized"}, status=403)

#_____________________________________________________________________________________________#

#Cart_View

def add_to_cart(request):
    if request.method == 'POST':
        # Get the product id and size from the form submission
        product_id = request.POST.get('product_id')
        size = request.POST.get('size')
        
        # Get the product object based on the provided product_id
        product = get_object_or_404(Product, id=product_id)

        # Get or create the user's cart (assuming the user is logged in)
        cart, created = Cart.objects.get_or_create(user=request.user)

        # Check if the product with the selected size already exists in the cart
        cart_item, created = CartItem.objects.get_or_create(
            cart=cart,
            product=product,
            size=size,
            defaults={'quantity': 1}  # Add 1 item if the cart item is new
        )

        if not created:
            # If the cart item already exists, just increment the quantity
            cart_item.quantity += 1
            cart_item.save()

        # Redirect to the cart page or any page you prefer after adding to cart
        return redirect('view_cart')  # Replace 'cart' with your actual cart URL name if needed

    # If the method is not POST, just redirect to the home or products page
    return redirect('home') 


def remove_from_cart(request):
    if request.method == 'POST':
        # Get the product ID from the form submission
        product_id = request.POST.get('product_id')

        # Get the user's cart
        cart = get_object_or_404(Cart, user=request.user)

        # Find the cart item and delete it
        cart_item = CartItem.objects.filter(cart=cart, product_id=product_id).first()

        if cart_item:
            cart_item.delete()

        # Redirect back to the cart page
        return redirect('view_cart')

    return redirect('home')

def product_detail(request, product_id):
    product = get_object_or_404(Product, id=product_id)
    sizes = range(5, 13)  # Sizes from 5 to 12
    return render(request, 'user/product_detail.html', {'product': product, 'sizes': sizes})

def view_cart(request):
    # Get the user's cart
    try:
        cart = Cart.objects.get(user=request.user)
    except Cart.DoesNotExist:
        cart = None

    # If the cart exists, fetch the cart items
    if cart:
        cart_items = cart.cart_items.all()
    else:
        cart_items = []

    total_price = sum(item.product.price * item.quantity for item in cart_items) 

    # Render the cart page with the cart items
    return render(request, 'user/cart.html', {'cart_items': cart_items, 'total_price': total_price})

def view_categories(request):
    categories = ProductCategory.objects.all()  # Fetch all categories
    return render(request, 'user/categories.html', {'categories': categories})

def view_products_by_category(request, category_id):
    sizes = [5, 6, 7, 8, 9, 10, 11, 12] 
    category = get_object_or_404(ProductCategory, id=category_id)  # Get the category
    products = Product.objects.filter(category=category)  # Get products by category
    return render(request, 'user/products_by_category.html', {'category': category, 'products': products,'sizes':sizes})


def search_view(request):
    query = request.GET.get('q', '')  # Get the search query from the GET request

    # Initialize empty result sets
    products = Product.objects.none()
    sizes = [5, 6, 7, 8, 9, 10, 11, 12] 

    if query:
        # Search for products by name or description
        products = Product.objects.filter(
            Q(name__icontains=query) | Q(description__icontains=query)
        )
       
    return render(request, 'user/search_results.html', {'products': products, 'query': query,'sizes':sizes})

#______________________________________________________________________________________________________________#
# Register Supplier
def register_supplier(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        email = request.POST.get('email')
        password = request.POST.get('password')
        confirm_password = request.POST.get('confirm_password')  # Get confirm password field
        company_name = request.POST.get('company_name')
        contact_info = request.POST.get('contact_info')

        if password != confirm_password:
            messages.error(request, 'Passwords do not match.')
            return redirect('register_supplier')

        if User.objects.filter(username=username).exists():
            messages.error(request, 'Username already exists.')
            return redirect('register_supplier')

        supplier_user = User.objects.create(
            username=username,
            email=email,
            password=make_password(password),
            role='supplier',
        )
        Supplier.objects.create(user=supplier_user, company_name=company_name, contact_info=contact_info)
        messages.success(request, 'Supplier registered successfully!')
        return redirect('login')

    return render(request, 'supplier/register_supplier.html')

def register_user(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        email = request.POST.get('email')
        password = request.POST.get('password')
        confirm_password = request.POST.get('confirm_password')  # Get confirm password field

        if password != confirm_password:
            messages.error(request, 'Passwords do not match.')
            return redirect('register_user')

        if User.objects.filter(username=username).exists():
            messages.error(request, 'Username already exists.')
            return redirect('register_user')

        User.objects.create(
            username=username,
            email=email,
            password=make_password(password),
            role='user',
            is_approved=True
        )
        messages.success(request, 'User registered successfully!')
        return redirect('login')

    return render(request, 'user/register_user.html')

# Login for All Roles
def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Authenticate the user using the custom User model
        user = authenticate(request, username=username, password=password)

        if user is not None:
            if user.is_superuser:
                    return redirect('/admin_dashboard/')
            # Check if the user is approved
            elif user.is_approved:
                login(request, user)
                messages.success(request, 'Login successful!')

                # Debugging output to check if user and role are correct
                print(f"User: {user}, Role: {user.role}")

                # Redirect based on the user's role
                if user.role == 'supplier':
                    return redirect('/supplier_dashboard/')
                elif user.role == 'user':
                    return redirect('/user_dashboard/')
                else:
                    messages.error(request, 'Invalid role assigned to the user.')
                    return redirect('login')
            else:
                messages.error(request, 'Your account is not approved yet. Please contact the admin.')
                return redirect('login')
        else:
            # If authentication fails, show an error message
            messages.error(request, 'Invalid username or password.')
            return redirect('login')

    return render(request, 'login.html')

#___________________________________________________________________________________________________#
def add_category(request):
    if request.method == 'POST':
        category_name = request.POST.get('category_name')
        category_description = request.POST.get('category_description', '')

        if category_name:
            try:

                # Create the new category for the logged-in supplier
                ProductCategory.objects.create(
                    name=category_name,
                    description=category_description,
                )
                messages.success(request, 'Category added successfully!')
                return redirect('view_all_categories')  # Redirect to the page displaying all categories
            except Supplier.DoesNotExist:
                # Handle the case where the user is not associated with a supplier
                messages.error(request, 'You must be a supplier to add categories.')
                return redirect('supplier_dashboard')  # Redirect to the supplier dashboard or another appropriate page
        else:
            messages.error(request, 'Category name is required!')

    return render(request, 'admin/add_category.html')

# Update Category for Supplier
def update_category(request, category_id):
    category = get_object_or_404(ProductCategory, id=category_id)
    if request.method == 'POST':
        category.name = request.POST.get('category_name', category.name)
        category.description = request.POST.get('category_description', category.description)
        category.save()
        messages.success(request, 'Category updated successfully!')
        return redirect('view_all_categories')  # Redirect to the page displaying all categories
    return render(request, 'admin/update_category.html', {'category': category})

# Delete Category for Supplier
def delete_category(request, category_id):
    category = get_object_or_404(ProductCategory, id=category_id)
    category.delete()
    messages.success(request, 'Category deleted successfully!')
    return redirect('view_all_categories')  # Redirect to the page displaying all categories

def add_product(request):
    if request.method == 'POST':
        product_name = request.POST.get('product_name')
        product_category_id = request.POST.get('category')
        product_description = request.POST.get('description', '')
        product_price = request.POST.get('price')
        product_stock = request.POST.get('stock')
        product_image = request.FILES.get('image')  # Handle image upload

        if product_name and product_price and product_stock:
            category = get_object_or_404(ProductCategory, id=product_category_id)
            Product.objects.create(
                name=product_name,
                category=category,
                description=product_description,
                price=product_price,
                stock=product_stock,
                image_url=product_image,  # Save the uploaded image
            )
            messages.success(request, 'Product added successfully!')
            return redirect('view_all_products')  # Redirect to the page displaying all products
        else:
            messages.error(request, 'Product name, price, and stock are required!')

    categories = ProductCategory.objects.all()
    return render(request, 'admin/add_product.html', {'categories': categories})

def update_product(request, product_id):
    product = get_object_or_404(Product, id=product_id)
    if request.method == 'POST':
        product.name = request.POST.get('product_name', product.name)
        product.category = get_object_or_404(ProductCategory, id=request.POST.get('category', product.category.id))
        product.description = request.POST.get('description', product.description)
        product.price = request.POST.get('price', product.price)
        product.stock = request.POST.get('stock', product.stock)

        # Handle image update
        if 'image' in request.FILES:  # Check if a new image was uploaded
            product.image_url = request.FILES['image']

        product.save()
        messages.success(request, 'Product updated successfully!')
        return redirect('view_all_products')  # Redirect to the page displaying all products

    categories = ProductCategory.objects.all()
    return render(request, 'admin/update_product.html', {'product': product, 'categories': categories})

# Delete Product for Supplier
def delete_product(request, product_id):
    product = get_object_or_404(Product, id=product_id)
    product.delete()
    messages.success(request, 'Product deleted successfully!')
    return redirect('view_all_products')  # Redirect to the page displaying all products

def view_all_categories(request):
    # Filter categories by the logged-in supplier
    categories = ProductCategory.objects.all()
    return render(request, 'admin/view_all_categories.html', {'categories': categories})

# View all products added by the supplier
def view_all_products(request):

    products = Product.objects.all()
    return render(request, 'admin/view_all_products.html', {'products': products})


def view_products(request):
    products = Product.objects.all()
    return render(request, 'supplier/view_products.html', {'products': products})


def user_logout(request):
    # Log the user out
    logout(request)
    # Redirect to a page after logging out (usually a login page or home page)
    return redirect('index')


def list_update_requests(request):
    update_requests = ProductUpdateRequest.objects.all().order_by('-requested_at')
    return render(request, 'admin/update_list.html', {'update_requests': update_requests})

def create_update_request(request, product_id):
    product = get_object_or_404(Product, id=product_id)
    supplier = get_object_or_404(Supplier, user=request.user)

    # Check if a pending request already exists
    existing_request = ProductUpdateRequest.objects.filter(product=product, supplier=supplier, is_approved=False)
    dup_req = ProductUpdateRequest.objects.filter(product=product, supplier=supplier, is_approved=True)
    if existing_request.exists():
        messages.warning(request, "An update request for this product is already pending.")
        return redirect('view_products')

    if dup_req.exists():
        messages.warning(request, "You already got update access.")
        return redirect('view_products')
    
    # Create the request
    ProductUpdateRequest.objects.create(product=product, supplier=supplier)
    messages.success(request, "Product update request submitted successfully.")
    return redirect('view_products')


def approve_update_request(request, request_id):
    update_request = get_object_or_404(ProductUpdateRequest, id=request_id)
    update_request.is_approved = True
    update_request.approved_at = timezone.now()
    update_request.save()

    messages.success(request, f"Update request for '{update_request.product.name}' approved.")
    return redirect('list_update_requests')

def update_product_supplier(request, product_id):
    product = get_object_or_404(Product, id=product_id)
    supplier = get_object_or_404(Supplier, user=request.user)

    # 🔒 Check if there’s an approved update request for this product by the supplier
    approved_request_exists = ProductUpdateRequest.objects.filter(
        product=product, supplier=supplier, is_approved=True
    ).exists()

    if not approved_request_exists:
        messages.error(request, "You cannot update this product. An approved update request is required.")
        return redirect('view_products')

    if request.method == 'POST':
        product.name = request.POST.get('product_name', product.name)
        product.category = get_object_or_404(ProductCategory, id=request.POST.get('category', product.category.id))
        product.description = request.POST.get('description', product.description)
        product.price = request.POST.get('price', product.price)
        product.stock = request.POST.get('stock', product.stock)

        # Handle image update
        if 'image' in request.FILES:
            product.image_url = request.FILES['image']

        product.save()

        # 🚫 Optional: After update, remove the approved request so the supplier can't reuse it
        ProductUpdateRequest.objects.filter(product=product, supplier=supplier, is_approved=True).delete()

        messages.success(request, 'Product updated successfully!')
        return redirect('view_products')

    categories = ProductCategory.objects.all()
    return render(request, 'supplier/update_product.html', {'product': product, 'categories': categories})


# Initialize Razorpay client
razorpay_client = razorpay.Client(auth=(settings.RAZORPAY_KEY_ID, settings.RAZORPAY_KEY_SECRET))
def make_payment(request):
    try:
        cart = Cart.objects.get(user=request.user)
        cart_items = cart.cart_items.all()
        total_price = sum(item.product.price * item.quantity for item in cart_items)

        if total_price == 0:
            return redirect('view_cart')

        # Initialize Razorpay client
        client = razorpay.Client(auth=(settings.RAZORPAY_KEY_ID, settings.RAZORPAY_KEY_SECRET))
        order_data = {
            'amount': int(total_price * 100),  # Convert to paisa
            'currency': 'INR',
            'payment_capture': '1'
        }
        order = client.order.create(order_data)

        # ✅ Fix: Do NOT pass 'order' as an argument
        payment = Payment.objects.create(
            user=request.user,
            amount=total_price,
            razorpay_order_id=order['id'],  # Use order['id'], NOT 'order'
            status="Pending"
        )

        context = {
            'order_id': order['id'],
            'amount': total_price,
            'razorpay_key': settings.RAZORPAY_KEY_ID,
        }
        return render(request, 'user/make_payment.html', context)

    except Cart.DoesNotExist:
        return redirect('view_cart')



def payment_success(request):
    """ Handles payment success """
    if request.method == "POST":
        razorpay_order_id = request.POST.get("razorpay_order_id")
        razorpay_payment_id = request.POST.get("razorpay_payment_id")
        razorpay_signature = request.POST.get("razorpay_signature")

        try:
            payment = Payment.objects.get(razorpay_order_id=razorpay_order_id)

            # Verify payment
            params_dict = {
                "razorpay_order_id": razorpay_order_id,
                "razorpay_payment_id": razorpay_payment_id,
                "razorpay_signature": razorpay_signature,
            }

            try:
                razorpay_client.utility.verify_payment_signature(params_dict)
                payment.status = "completed"
                payment.razorpay_payment_id = razorpay_payment_id
                payment.razorpay_signature = razorpay_signature
                payment.save()

                cart_items = CartItem.objects.filter(cart=payment.user.cart)

                # Create Order and associate it with Payment
                order = Order.objects.create(
                user=request.user,
                total_amount=payment.amount,
                payment=payment
                )
                order.products.set([item.product for item in cart_items])  


                # Reduce product stock
                for cart_item in cart_items:
                    product = cart_item.product
                    if product.stock >= cart_item.quantity:
                        product.stock -= cart_item.quantity  # Deduct quantity purchased
                        product.save()
                    else:
                        messages.error(request, f"Not enough stock for {product.name}.")
                        return redirect("view_cart")
                
                # Clear the cart
                cart_items.delete()

                messages.success(request, "Payment successful! Your order has been placed.")
                return redirect("view_cart")

            except razorpay.errors.SignatureVerificationError:
                messages.error(request, "Payment verification failed.")
                return redirect("payment_failure")

        except Payment.DoesNotExist:
            messages.error(request, "Invalid payment attempt.")
            return redirect("view_cart")

    return redirect("view_cart")



def payment_failure(request):
    """ Handles payment failure """
    messages.error(request, "Payment failed. Please try again.")
    return redirect("view_cart")

def my_orders(request):
    """ View to display the current logged-in user's orders with products """
    orders = Order.objects.filter(user=request.user).prefetch_related('products')  # Optimize query
    
    return render(request, 'user/my_orders.html', {'orders': orders})

def download_invoice(request, order_id):
    """ Generate a PDF invoice dynamically for an order """
    order = get_object_or_404(Order, id=order_id, user=request.user)

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="invoice_{order.id}.pdf"'

    p = canvas.Canvas(response, pagesize=letter)
    width, height = letter

    # Invoice Title
    p.setFont("Helvetica-Bold", 16)
    p.drawString(200, height - 50, "Invoice")

    # Order Details
    p.setFont("Helvetica", 12)
    p.drawString(50, height - 100, f"Order ID: {order.id}")
    p.drawString(50, height - 120, f"Order Date: {order.created_at.strftime('%d %b %Y, %H:%M')}")
    p.drawString(50, height - 140, f"Customer: {order.user.username}")
    p.drawString(50, height - 160, f"Email: {order.user.email}")

    # Products
    p.setFont("Helvetica-Bold", 12)
    p.drawString(50, height - 200, "Products:")

    y_position = height - 220
    p.setFont("Helvetica", 12)

    for product in order.products.all():
        p.drawString(50, y_position, f"{product.name} - ${product.price}")
        y_position -= 20  # Move down for next product

    # Total Amount
    p.setFont("Helvetica-Bold", 12)
    p.drawString(50, y_position - 20, f"Total Amount: ${order.total_amount}")

    # Finish PDF
    p.showPage()
    p.save()

    return response
